from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # slide backgrounds
ACCENT_BLUE = RGBColor(0x1E, 0x6F, 0xBF)   # headings / highlights
ACCENT_TEAL = RGBColor(0x17, 0xA2, 0xB8)   # accent dots / icons
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY    = RGBColor(0x8A, 0x9B, 0xB0)
GREEN       = RGBColor(0x28, 0xA7, 0x45)
ORANGE      = RGBColor(0xFD, 0x7E, 0x14)
RED         = RGBColor(0xDC, 0x35, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    return slide

def fill_bg(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=WHITE, bold=False,
                  leading_color=None, leading_char="●"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        if leading_color and line:
            run0 = p.add_run()
            run0.text = leading_char + "  "
            run0.font.size = Pt(font_size)
            run0.font.color.rgb = leading_color
            run0.font.bold = True
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def slide_header(slide, title, subtitle=None):
    """Top accent bar + title."""
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT_BLUE)
    add_text(slide, title, 0.5, 0.18, 12.3, 0.7,
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12.3, 0.45,
                 font_size=14, color=ACCENT_TEAL)
    add_rect(slide, 0.5, 1.35, 1.2, 0.04, ACCENT_TEAL)


def section_card(slide, left, top, width, height, title, lines,
                 card_color=RGBColor(0x16, 0x2B, 0x45),
                 title_color=ACCENT_TEAL, text_color=WHITE,
                 font_size=12):
    add_rect(slide, left, top, width, height, card_color)
    add_text(slide, title, left+0.15, top+0.12, width-0.3, 0.4,
             font_size=13, bold=True, color=title_color)
    add_multiline(slide, lines, left+0.15, top+0.52, width-0.3,
                  height-0.65, font_size=font_size, color=text_color,
                  leading_color=ACCENT_TEAL)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)

# large decorative rect
add_rect(s, 0, 0, 5.2, 7.5, ACCENT_BLUE)
add_rect(s, 5.2, 0, 8.13, 7.5, DARK_NAVY)

# badge
add_rect(s, 5.5, 1.5, 3, 0.55, ACCENT_TEAL)
add_text(s, "PILOT PROJECT  ·  APPROVED v1.1", 5.5, 1.56, 3, 0.45,
         font_size=11, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

add_text(s, "Authentication\nSystem", 5.5, 2.2, 7.3, 2.0,
         font_size=46, bold=True, color=WHITE)
add_text(s, "Self-Hosted · Multi-Tenant · Standards-Compliant",
         5.5, 4.25, 7.3, 0.6, font_size=16, color=ACCENT_TEAL)
add_text(s, "Business Overview  ·  March 2026",
         5.5, 4.9, 7.3, 0.5, font_size=13, color=MID_GRAY)

# left panel text
add_text(s, "Digital Worker", 0.35, 0.4, 4.5, 0.6,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Eliminating vendor lock-in\nand security gaps across\nall our products.",
         0.35, 1.4, 4.5, 2.0, font_size=18, color=WHITE)
add_text(s, "PRD v1.1 · Status: ✅ Approved",
         0.35, 6.8, 4.5, 0.45, font_size=11, color=RGBColor(0xA0,0xC4,0xFF))


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "The Problem We're Solving",
             "Why build our own authentication platform?")

# Three pain columns
cols = [
    ("End Users", ["Separate login per product", "No Single Sign-On", "Poor password reset UX"]),
    ("Developers", ["Re-implement auth in every service", "Inconsistent security guarantees", "Weeks to integrate"]),
    ("Business", ["Vendor lock-in (Auth0/Okta)", "Unpredictable per-user pricing", "No control over identity data"]),
]
x = 0.5
for title, items in cols:
    section_card(s, x, 1.55, 3.9, 2.3, title, items, font_size=13)
    x += 4.1

# Bottom row
section_card(s, 0.5, 3.95, 5.9, 2.2,
             "Security & Compliance Teams",
             ["No central audit log", "Cannot demonstrate GDPR / SOC2 compliance",
              "Difficulty revoking sessions across services"],
             font_size=13)
section_card(s, 6.5, 3.95, 6.3, 2.2,
             "Tenant Administrators",
             ["No visibility into who is logged in", "Cannot revoke sessions",
              "No exportable audit logs for their org"],
             font_size=13)

add_text(s, "Current vendor cost scales unpredictably at volume  ·  Sensitive identity data lives outside our control",
         0.5, 6.3, 12.3, 0.5, font_size=11, color=MID_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Our Solution", "A self-hosted, standards-compliant authentication platform")

# Centre pill
add_rect(s, 4.2, 1.6, 4.9, 1.4, ACCENT_BLUE)
add_text(s, "Authentication Platform", 4.2, 1.72, 4.9, 0.55,
         font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Self-hosted  ·  Multi-tenant  ·  API-first",
         4.2, 2.2, 4.9, 0.45, font_size=11, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

cards = [
    (0.3, 1.5, 3.5, "Replaces", ["Auth0, Okta, Cognito", "→ Zero per-MAU cost", "→ Full data ownership"]),
    (9.4, 1.5, 3.5, "Serves", ["End Users via SSO", "Developers via OAuth 2.0", "Admins via Tenant APIs"]),
    (0.3, 3.5, 3.5, "Built on", ["OAuth 2.0 / OIDC standards", "JWT + Argon2id security", "PostgreSQL + Redis"]),
    (9.4, 3.5, 3.5, "Enables", ["SOC2 & GDPR compliance", "Multi-org isolation", "Integration in < 2 days"]),
]
for (lft, tp, wd, ttl, items) in cards:
    section_card(s, lft, tp, wd, 1.8, ttl, items, font_size=12)

# connecting lines (visual only via thin rects)
add_rect(s, 3.8, 2.25, 0.4, 0.05, ACCENT_TEAL)
add_rect(s, 9.1, 2.25, 0.35, 0.05, ACCENT_TEAL)
add_rect(s, 3.8, 4.35, 0.4, 0.05, ACCENT_TEAL)
add_rect(s, 9.1, 4.35, 0.35, 0.05, ACCENT_TEAL)

add_text(s, "API-only design · No vendor lock-in · Identity data stays on our infrastructure",
         0.5, 6.75, 12.3, 0.45, font_size=12, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — BUSINESS GOALS & SUCCESS METRICS
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Business Goals & Success Metrics",
             "Measurable outcomes we commit to delivering within 6 months of launch")

# Goals (left)
section_card(s, 0.5, 1.55, 5.8, 2.1, "Business Goals",
             ["Eliminate per-MAU licensing costs from third-party auth vendors",
              "Reduce new service auth integration from weeks → ≤ 2 days",
              "Achieve compliance posture for SOC2 Type II and GDPR audits"],
             font_size=13)

# Metrics (right) — table style
metrics = [
    ("New service integration time",  "≤ 2 business days"),
    ("Login latency (p95)",           "< 300 ms"),
    ("Auth error rate",               "< 0.1%"),
    ("MFA adoption — admin users",    "≥ 60%"),
    ("Security incidents",            "0"),
    ("Audit log completeness",        "100% of events"),
    ("Tenant onboarding",             "< 30 min self-service"),
    ("Developer satisfaction",        "≥ 4.2 / 5"),
]
add_rect(s, 6.5, 1.55, 6.3, 2.0, RGBColor(0x16,0x2B,0x45))
add_text(s, "Success Metrics  (6 months post-launch)", 6.65, 1.62, 6.0, 0.45,
         font_size=13, bold=True, color=ACCENT_TEAL)
row_y = 2.05
for metric, target in metrics:
    add_text(s, metric, 6.65, row_y, 4.1, 0.28, font_size=10.5, color=WHITE)
    add_text(s, target, 10.8, row_y, 1.8, 0.28, font_size=10.5, bold=True, color=GREEN)
    row_y += 0.28

# Key win callouts
kw = [
    ("Cost Control",    "No surprise per-user bills\nas we scale"),
    ("Speed",           "New services live in days,\nnot weeks"),
    ("Compliance",      "SOC2 + GDPR evidence\nfrom day one"),
    ("Ownership",       "Identity data fully\nunder our control"),
]
x = 0.5
for title, body in kw:
    add_rect(s, x, 3.85, 2.95, 1.5, RGBColor(0x0A,0x3D,0x62))
    add_text(s, title, x+0.15, 3.95, 2.65, 0.45, font_size=13, bold=True, color=ACCENT_TEAL)
    add_text(s, body,  x+0.15, 4.38, 2.65, 0.9,  font_size=12, color=WHITE)
    x += 3.1


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — WHO BENEFITS  (Personas)
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Who Benefits From This System",
             "Four personas — each with distinct needs and measurable success criteria")

personas = [
    ("Maya\nEnd User", ACCENT_BLUE,
     ["Fast, frictionless login",
      "Social login (Google)",
      "Self-service password reset",
      "Optional MFA for security",
      "Clear error messages"]),
    ("Carlos\nTenant Admin", RGBColor(0x0A,0x3D,0x62),
     ["Invite / disable / remove users",
      "Assign roles & permissions",
      "Access audit logs",
      "Enforce MFA for his org",
      "Onboard new hire in < 5 min"]),
    ("Priya\nDeveloper", RGBColor(0x14,0x4E,0x3C),
     ["Standard OAuth 2.0 endpoints",
      "JWT tokens with role claims",
      "Clear error codes & docs",
      "Sandbox tenant for testing",
      "Integrate in < half a day"]),
    ("Jordan\nPlatform Admin", RGBColor(0x4A,0x1A,0x5E),
     ["Provision tenants via API",
      "Cross-tenant audit visibility",
      "System health monitoring",
      "Incident response tools",
      "Provision new tenant in < 10 min"]),
]
x = 0.4
colors = [ACCENT_BLUE, RGBColor(0x0A,0x3D,0x62), RGBColor(0x14,0x4E,0x3C), RGBColor(0x4A,0x1A,0x5E)]
for (name, col, items), card_col in zip(personas, colors):
    add_rect(s, x, 1.55, 3.0, 4.5, card_col)
    add_text(s, name, x+0.15, 1.65, 2.7, 0.75,
             font_size=16, bold=True, color=WHITE)
    add_rect(s, x+0.15, 2.4, 2.7, 0.04, ACCENT_TEAL)
    add_multiline(s, items, x+0.15, 2.55, 2.7, 3.3,
                  font_size=12, color=WHITE, leading_color=ACCENT_TEAL)
    x += 3.15

add_text(s, "The auth platform serves all four simultaneously through a single, consistent API.",
         0.5, 6.25, 12.3, 0.5, font_size=12, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — FEATURE SCOPE  (MoSCoW)
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Feature Scope — What We're Building",
             "18 Must-Have features form the MVP delivered over 9 sprints")

must = [
    "User registration & email verification",
    "Login (password) + password reset",
    "Secure session management (JWT + refresh token rotation)",
    "Multi-tenant data isolation (schema-per-tenant)",
    "Tenant provisioning API",
    "Role-Based Access Control (RBAC)",
    "OAuth 2.0 — Authorization Code + PKCE",
    "OAuth 2.0 — Client Credentials (M2M / service-to-service)",
    "Social Login via Google",
    "Rate limiting & brute-force protection",
    "Comprehensive audit log (append-only)",
    "HTTPS enforcement + secure response headers",
    "Admin user management API",
    "Token introspection & JWKS endpoints",
]
should = [
    "MFA via TOTP (authenticator app)",
    "Per-tenant MFA enforcement",
    "Microsoft / Entra ID social login",
    "Fine-grained permission model",
    "Self-service user profile",
    "Webhook notifications on auth events",
    "Session activity API",
    "Developer documentation site",
]
wont = [
    "SAML 2.0 (enterprise phase)",
    "Hosted login UI pages",
    "Multi-region active-active",
    "Native mobile SDKs",
    "LDAP / Kerberos",
    "Cross-tenant user identities",
    "Billing / subscription management",
]

# Must Have card
add_rect(s, 0.4, 1.55, 5.0, 5.3, RGBColor(0x0A,0x2B,0x18))
add_text(s, "✅  Must Have  (MVP — 18 features)", 0.55, 1.62, 4.7, 0.45,
         font_size=13, bold=True, color=GREEN)
add_multiline(s, must, 0.55, 2.1, 4.7, 4.6,
              font_size=10.5, color=WHITE, leading_color=GREEN)

# Should Have card
add_rect(s, 5.55, 1.55, 3.9, 5.3, RGBColor(0x2B,0x1E,0x06))
add_text(s, "🔶  Should Have", 5.7, 1.62, 3.6, 0.45,
         font_size=13, bold=True, color=ORANGE)
add_multiline(s, should, 5.7, 2.1, 3.6, 4.6,
              font_size=10.5, color=WHITE, leading_color=ORANGE)

# Won't Do card
add_rect(s, 9.6, 1.55, 3.3, 5.3, RGBColor(0x2B,0x0A,0x0A))
add_text(s, "🚫  Won't Do (This Pilot)", 9.75, 1.62, 3.0, 0.45,
         font_size=13, bold=True, color=RED)
add_multiline(s, wont, 9.75, 2.1, 3.0, 4.6,
              font_size=10.5, color=WHITE, leading_color=RED)

add_text(s, "Scope is fixed. Any additions require a formal change request with re-estimation.",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — ARCHITECTURE OVERVIEW (Business-friendly)
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "How It Works — Architecture at a Glance",
             "Key design decisions that protect our data and enable fast scaling")

decisions = [
    ("Schema-per-Tenant\nIsolation",
     "Each organisation's data lives in its own database schema.\nCross-tenant data leakage is architecturally impossible.",
     ACCENT_BLUE),
    ("JWT + Refresh\nToken Strategy",
     "Short-lived access tokens (15 min) + long-lived opaque refresh tokens.\nResource servers validate locally — no central bottleneck.",
     RGBColor(0x0A,0x3D,0x62)),
    ("Secrets Manager\n(Vault)",
     "All signing keys and credentials stored in a dedicated vault.\nNever in code, config files, or environment variables.",
     RGBColor(0x14,0x4E,0x3C)),
    ("API-Only\n(No Hosted UI)",
     "Clean, testable API surface. Consuming apps build their own login UI.\nReduces scope by 2–4 sprints; forces good API-first design.",
     RGBColor(0x4A,0x1A,0x5E)),
    ("OAuth 2.0 +\nPKCE Standard",
     "Industry-standard protocol. Every integration is immediately\nportable and auditor-recognisable.",
     RGBColor(0x5A,0x2D,0x00)),
    ("1-Year Audit Log\n+ Cold Archive",
     "Every auth event logged with user, IP, timestamp, outcome.\nSatisfies SOC2 Type II evidence requirements.",
     RGBColor(0x1A,0x3A,0x5C)),
]
x, y = 0.4, 1.55
for i, (title, desc, col) in enumerate(decisions):
    if i == 3:
        x, y = 0.4, 4.25
    add_rect(s, x, y, 4.1, 1.55, col)
    add_text(s, title, x+0.15, y+0.1, 3.8, 0.6, font_size=13, bold=True, color=WHITE)
    add_text(s, desc,  x+0.15, y+0.72, 3.8, 0.75, font_size=10.5, color=LIGHT_GRAY)
    x += 4.3

add_text(s, "All architectural decisions are finalised (ADR-001 through ADR-009) and approved for development.",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — DELIVERY TIMELINE
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Delivery Timeline — 9 Sprints · ~20 Weeks",
             "3 backend engineers · 1 QA · 0.5 DevOps · Velocity: 28–32 pts/sprint")

sprints = [
    ("S0", "Wk 1–2",  "Architecture\nSpike",          "—",   ACCENT_TEAL,   "ERD · CI/CD · Staging env"),
    ("S1", "Wk 3–4",  "Core Auth\nFoundation",         "21",  GREEN,         "Register · Login · Email verify · Password reset"),
    ("S2", "Wk 5–6",  "Sessions,\nSecurity, Audit",    "20",  GREEN,         "Session refresh · Logout · Rate limiting · Audit log"),
    ("S3", "Wk 7–8",  "Multi-Tenancy\nFoundation",     "23",  ORANGE,        "Tenant provisioning · Schema isolation · Migration runner  ⚠️ Highest risk"),
    ("S4", "Wk 9–10", "RBAC &\nAdmin API",             "15",  GREEN,         "Role assignment · JWT claims · Admin user management"),
    ("S5", "Wk 11–12","OAuth 2.0\nAuth Server",        "11",  GREEN,         "Auth Code + PKCE · Client registration"),
    ("S6", "Wk 13–14","M2M + Google\nSocial Login",    "13",  GREEN,         "Service-to-service tokens · Google OAuth"),
    ("S7", "Wk 15–16","MFA + User\nProfile",           "16",  GREEN,         "TOTP MFA · MFA enforcement · Self-service profile"),
    ("S8", "Wk 17–18","Hardening &\nCompliance",       "25",  ORANGE,        "Pentest remediation · GDPR · Perf testing · Docs"),
    ("S9", "Wk 19–20","Launch\nPreparation",           "15",  ACCENT_TEAL,   "E2E testing · Go/no-go · Prod smoke tests"),
]

bar_x_start = 1.7
bar_total_w = 11.0
total_pts = 159

y = 1.45
for snum, weeks, theme, pts, col, notes in sprints:
    # sprint label
    add_text(s, snum,   0.15, y, 0.6, 0.42, font_size=11, bold=True, color=col)
    add_text(s, weeks,  0.75, y, 0.9, 0.42, font_size=9,  color=MID_GRAY)
    # bar
    bar_w = (int(pts) / total_pts * bar_total_w) if pts != "—" else 0.4
    add_rect(s, bar_x_start, y+0.08, bar_w, 0.32, col)
    # theme inside bar (or beside if too narrow)
    if bar_w > 1.5:
        add_text(s, theme, bar_x_start+0.1, y+0.06, bar_w-0.15, 0.38,
                 font_size=9, bold=True, color=DARK_NAVY)
    else:
        add_text(s, theme, bar_x_start+bar_w+0.05, y+0.06, 2.0, 0.38,
                 font_size=9, bold=True, color=WHITE)
    # pts badge
    if pts != "—":
        add_text(s, f"{pts} pts", bar_x_start+bar_w+0.1, y+0.1, 0.7, 0.28,
                 font_size=9, bold=True, color=col)
    # notes
    note_x = bar_x_start + bar_w + (0.85 if pts != "—" else 0.5)
    add_text(s, notes, note_x, y+0.08, 13.33-note_x-0.3, 0.38,
             font_size=8.5, color=MID_GRAY)
    y += 0.52

add_text(s, "Total: ~159 story points  ·  ~20 weeks  ·  Team: 3 BE engineers, 1 QA, 0.5 DevOps",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — RISKS
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Key Risks & Mitigations",
             "Identified and planned for — not surprises")

risks = [
    (ORANGE, "Medium / High",  "Migration runner complexity (Sprint 3)",
             "Sprint 0 spike must deliver a working POC, not just a design. Non-negotiable gate."),
    (RED,    "High / Medium",  "OAuth 2.0 edge-case overrun (Sprint 5)",
             "Use a battle-tested OAuth library. Strict sprint timebox enforced."),
    (RED,    "Low / Critical", "Cross-tenant data leak",
             "Automated isolation test suite (US-08b) is non-negotiable Definition of Done for every data story."),
    (ORANGE, "Medium / High",  "Pentest uncovers Critical findings",
             "Book penetration tester in Sprint 0 for Sprint 6/7 window. Do not wait for Sprint 8."),
    (ORANGE, "High / Medium",  "Scope creep — SAML / other enterprise features",
             "SAML is explicitly Out of Scope (OS-01). Any addition requires formal change request + re-estimation."),
    (ORANGE, "Medium / High",  "Refresh token security edge cases",
             "Designate auth security champion. Mandatory RFC 6749/7009/7519/7636 review before Sprint 2."),
]

y = 1.55
for col, likelihood, risk, mitigation in risks:
    add_rect(s, 0.4, y, 0.08, 0.72, col)
    add_rect(s, 0.55, y, 12.3, 0.72, RGBColor(0x16,0x2B,0x45))
    add_text(s, f"[{likelihood}]  {risk}", 0.7, y+0.04, 8.5, 0.35,
             font_size=12, bold=True, color=WHITE)
    add_text(s, f"Mitigation: {mitigation}", 0.7, y+0.38, 11.8, 0.3,
             font_size=10.5, color=MID_GRAY, italic=True)
    y += 0.85

add_text(s, "Risk register reviewed at start of every sprint. Early escalation is expected.",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — COMPLIANCE & SECURITY
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Security & Compliance Commitments",
             "Built-in from day one — not bolted on at the end")

left_items = [
    ("Password Security",      "Argon2id hashing — the gold standard. MD5/SHA1 rejected."),
    ("Token Security",         "RS256-signed JWTs. 15-minute access token TTL. Refresh tokens rotated on every use."),
    ("Transport",              "TLS 1.3 preferred. HSTS enforced. No HTTP fallback in production."),
    ("Secrets Management",     "All keys & credentials in HashiCorp Vault. Never in code or config files."),
    ("Brute Force Protection", "Per-IP and per-user rate limiting backed by Redis. Configurable per tenant."),
    ("OAuth Security",         "PKCE S256 mandatory for all authorization code flows. state param required."),
]
right_items = [
    ("GDPR",      "Right-to-erasure API. PII deletion across tenant schema. Audit entries anonymised."),
    ("SOC2 CC6.2","All auth events logged: user, IP, timestamp, outcome. 100% completeness target."),
    ("SOC2 CC6.3","Access revocation takes effect within one access token TTL (15 minutes)."),
    ("Audit Log", "Append-only. 1-year hot retention + cold archive. Tenant-admin read-only access."),
    ("Penetration Test", "External pentest scheduled for Sprint 6/7 window. Findings remediated in Sprint 8."),
    ("OWASP Top 10", "Full review in Sprint 8. SQL injection and XSS prevented by design (parameterised queries, CSP headers)."),
]

y = 1.55
for title, desc in left_items:
    add_rect(s, 0.4, y, 5.9, 0.7, RGBColor(0x16,0x2B,0x45))
    add_text(s, title, 0.55, y+0.05, 5.6, 0.28, font_size=11, bold=True, color=ACCENT_TEAL)
    add_text(s, desc,  0.55, y+0.35, 5.6, 0.3,  font_size=10, color=WHITE)
    y += 0.78

y = 1.55
for title, desc in right_items:
    add_rect(s, 6.85, y, 6.1, 0.7, RGBColor(0x0A,0x2B,0x18))
    add_text(s, title, 7.0, y+0.05, 5.8, 0.28, font_size=11, bold=True, color=GREEN)
    add_text(s, desc,  7.0, y+0.35, 5.8, 0.3,  font_size=10, color=WHITE)
    y += 0.78

add_text(s, "SOC2 Type II evidence collection begins at launch · Compliance consultant engaged by Sprint 2",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=ACCENT_TEAL, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11 — NEXT STEPS
# ══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, DARK_NAVY)
slide_header(s, "Next Steps & Decisions Required",
             "Status: PRD v1.1 approved · Ready for development kickoff")

actions = [
    ("✅  DONE — No Action Needed",  GREEN,  [
        "PRD v1.1 finalised — all open questions resolved",
        "9 Architecture Decision Records (ADR-001 to ADR-009) locked",
        "Full product backlog estimated (~159 points)",
        "Sprint plan defined across 9 sprints / ~20 weeks",
    ]),
    ("🔶  In Progress — Confirm This Week", ORANGE, [
        "Sprint 0 kickoff — architecture spike tasks assigned",
        "Book external penetration tester (Sprint 6/7 window)",
        "Engage compliance consultant (target: before Sprint 2)",
        "Provision staging: PostgreSQL + Redis + secrets manager",
    ]),
    ("📋  Business Decisions Needed", ACCENT_TEAL, [
        "Confirm team staffing: 3 BE + 1 QA + 0.5 DevOps available?",
        "Sign off on default rate-limiting / lockout thresholds",
        "Confirm Google Cloud project for OAuth credentials (Sprint 6)",
        "Approve formal scope change process for any new requirements",
    ]),
]

y = 1.55
for header, col, items in actions:
    add_rect(s, 0.4, y, 12.4, 1.7, RGBColor(0x16,0x2B,0x45))
    add_rect(s, 0.4, y, 0.08, 1.7, col)
    add_text(s, header, 0.6, y+0.1, 11.9, 0.45, font_size=13, bold=True, color=col)
    add_multiline(s, items, 0.6, y+0.55, 11.9, 1.1,
                  font_size=12, color=WHITE, leading_color=col)
    y += 1.9

add_text(s, "Questions? Reach out to the Product Owner — PRD v1.1 is the single source of truth.",
         0.5, 6.95, 12.3, 0.4, font_size=11, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════
out_path = "/Users/kanatekhumnin/Project/digital-worker/docs/auth-system/auth-system-business-overview.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
